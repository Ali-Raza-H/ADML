import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from adml.adt.core import Slide, TextElement, ShapeElement, ListElement, ImageElement, FillType
from adml.renderers.base import BaseRenderer
from adml.utils.units import to_emu

class PowerPointRenderer(BaseRenderer):

    def get_file_extension(self) -> str:
        return "pptx"

    def render(self) -> bytes:
        prs = Presentation()
        prs.slide_width = to_emu(self.doc.width)
        prs.slide_height = to_emu(self.doc.height)
        blank_slide_layout = prs.slide_layouts[6]
        
        for slide in self.doc.slides:
            prs_slide = prs.slides.add_slide(blank_slide_layout)
            self._render_background(prs_slide, slide)
            self._render_slide(prs_slide, slide)
            
        out = BytesIO()
        prs.save(out)
        return out.getvalue()

    def _render_slide(self, prs_slide, slide: Slide):
        for el in slide.elements:
            if isinstance(el, TextElement):
                self._render_text(prs_slide, el, slide)
            elif isinstance(el, ShapeElement):
                self._render_shape(prs_slide, el, slide)
            elif isinstance(el, ListElement):
                self._render_list(prs_slide, el, slide)
            elif isinstance(el, ImageElement):
                self._render_image(prs_slide, el, slide)

    def _render_background(self, prs_slide, slide: Slide):
        if slide.background.type == FillType.SOLID and slide.background.colour:
            bg = prs_slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = slide.background.colour.to_pptx_colour()
        elif slide.background.type == FillType.GRADIENT:
            shape = prs_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, to_emu(slide.width), to_emu(slide.height)
            )
            shape.line.fill.background()
            shape.fill.solid()
            c = self.resolve_background_colour(slide.background)
            shape.fill.fore_color.rgb = c.to_pptx_colour()

    def _render_text(self, prs_slide, element: TextElement, slide: Slide):
        x_pt, y_pt = self.resolve_position(element, slide.width, slide.height)
        w_pt = element.width or slide.width * 0.8
        h_pt = element.height or element.font.size * 2.5
        
        txBox = prs_slide.shapes.add_textbox(to_emu(x_pt), to_emu(y_pt), to_emu(w_pt), to_emu(h_pt))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = element.content
        
        if element.font.align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif element.font.align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
            
        for run in p.runs:
            run.font.name = element.font.family
            run.font.size = Pt(element.font.size)
            if element.font.colour:
                run.font.color.rgb = element.font.colour.to_pptx_colour()
            if element.font.weight in ("bold", "extrabold", "black", "semibold"):
                run.font.bold = True
            if element.font.style == "italic":
                run.font.italic = True

    def _render_shape(self, prs_slide, element: ShapeElement, slide: Slide):
        x_pt, y_pt = self.resolve_position(element, slide.width, slide.height)
        w_pt = element.width or 100.0
        h_pt = element.height or 100.0
        
        shape_type = MSO_SHAPE.RECTANGLE
        if element.shape_type == "circle" or element.shape_type == "ellipse":
            shape_type = MSO_SHAPE.OVAL
        elif element.shape_type == "rectangle" and element.corner_radius > 0:
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        
        if element.shape_type == "line":
            shape = prs_slide.shapes.add_connector(MSO_SHAPE.LINE, to_emu(x_pt), to_emu(y_pt), to_emu(x_pt + w_pt), to_emu(y_pt + h_pt))
        else:
            shape = prs_slide.shapes.add_shape(shape_type, to_emu(x_pt), to_emu(y_pt), to_emu(w_pt), to_emu(h_pt))
            if element.shape_type == "rectangle" and element.corner_radius > 0:
                adj_val = min(element.corner_radius / min(w_pt, h_pt), 0.5) * 100000
                if len(shape.adjustments) > 0:
                    shape.adjustments[0] = int(adj_val)

        if element.fill.type == FillType.SOLID and element.fill.colour:
            shape.fill.solid()
            shape.fill.fore_color.rgb = element.fill.colour.to_pptx_colour()
        else:
            shape.fill.background()
            
        if element.stroke_width > 0 and element.stroke_colour:
            shape.line.width = Pt(element.stroke_width)
            shape.line.color.rgb = element.stroke_colour.to_pptx_colour()
        else:
            shape.line.fill.background()

    def _render_list(self, prs_slide, element: ListElement, slide: Slide):
        x_pt, y_pt = self.resolve_position(element, slide.width, slide.height)
        w_pt = element.width or slide.width * 0.8
        h_pt = element.height or (element.font.size * 2 * len(element.items))
        
        txBox = prs_slide.shapes.add_textbox(to_emu(x_pt), to_emu(y_pt), to_emu(w_pt), to_emu(h_pt))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(element.items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            
            if element.list_type == "bullet":
                p.text = item
                p.level = 1
            elif element.list_type == "numbered":
                p.text = f"{i+1}. {item}"
                
            p.space_after = Pt(element.item_spacing)
            
            for run in p.runs:
                run.font.name = element.font.family
                run.font.size = Pt(element.font.size)
                if element.font.colour:
                    run.font.color.rgb = element.font.colour.to_pptx_colour()
                if element.font.weight in ("bold", "extrabold", "black", "semibold"):
                    run.font.bold = True
                if element.font.style == "italic":
                    run.font.italic = True

    def _render_image(self, prs_slide, element: ImageElement, slide: Slide):
        x_pt, y_pt = self.resolve_position(element, slide.width, slide.height)
        w_pt = element.width or 200.0
        h_pt = element.height or 200.0
        
        if os.path.exists(element.src):
            prs_slide.shapes.add_picture(element.src, to_emu(x_pt), to_emu(y_pt), to_emu(w_pt), to_emu(h_pt))
        else:
            shape = prs_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, to_emu(x_pt), to_emu(y_pt), to_emu(w_pt), to_emu(h_pt))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shape.text = element.src
